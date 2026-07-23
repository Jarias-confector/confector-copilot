from dataclasses import dataclass, field
from pathlib import Path

from confector_core import DocumentKind

_EXTENSION_KIND = {
    ".pdf": DocumentKind.PDF,
    ".doc": DocumentKind.WORD,
    ".docx": DocumentKind.WORD,
    ".xls": DocumentKind.EXCEL,
    ".xlsx": DocumentKind.EXCEL,
    ".csv": DocumentKind.CSV,
    ".ppt": DocumentKind.PRESENTATION,
    ".pptx": DocumentKind.PRESENTATION,
    ".txt": DocumentKind.TEXT,
    ".md": DocumentKind.TEXT,
    ".mp3": DocumentKind.AUDIO,
    ".wav": DocumentKind.AUDIO,
    ".m4a": DocumentKind.AUDIO,
}

_MAX_EXTRACT_CHARS = 5000


@dataclass
class Extraction:
    """Salida estandarizada de cualquier extractor — mismo shape sin importar el tipo de archivo."""

    text: str | None
    metadata: dict = field(default_factory=dict)


def classify(filename: str) -> DocumentKind:
    return _EXTENSION_KIND.get(Path(filename).suffix.lower(), DocumentKind.OTHER)


def extract(kind: DocumentKind, content: bytes, filename: str = "") -> Extraction:
    """Extracción inicial + metadatos (10_MVP_ROADMAP Sprint 2/3)."""
    try:
        if kind == DocumentKind.PDF:
            return _extract_pdf(content)
        if kind == DocumentKind.WORD:
            return _extract_docx(content)
        if kind == DocumentKind.EXCEL:
            return _extract_xlsx(content)
        if kind == DocumentKind.CSV:
            return _extract_csv(content)
        if kind == DocumentKind.PRESENTATION:
            return _extract_pptx(content)
        if kind == DocumentKind.TEXT:
            return _extract_text(content)
        if kind == DocumentKind.AUDIO:
            return _extract_audio(content, Path(filename).suffix or ".wav")
    except Exception:
        return Extraction(text=None, metadata={"error": "no se pudo extraer el contenido"})
    return Extraction(text=None)


def _extract_pdf(content: bytes) -> Extraction:
    import fitz  # PyMuPDF

    with fitz.open(stream=content, filetype="pdf") as doc:
        text = "\n".join(page.get_text() for page in doc)
        pages = doc.page_count
    return Extraction(text=text[:_MAX_EXTRACT_CHARS], metadata={"pages": pages, "words": len(text.split())})


def _extract_docx(content: bytes) -> Extraction:
    import io

    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)
    return Extraction(
        text=text[:_MAX_EXTRACT_CHARS],
        metadata={"paragraphs": len(paragraphs), "words": len(text.split())},
    )


def _extract_xlsx(content: bytes) -> Extraction:
    import io

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    row_count = 0
    for sheet in wb.worksheets:
        lines.append(f"# {sheet.title}")
        for row in sheet.iter_rows(max_row=20, values_only=True):
            lines.append(", ".join(str(cell) for cell in row if cell is not None))
            row_count += 1
    return Extraction(
        text="\n".join(lines)[:_MAX_EXTRACT_CHARS],
        metadata={"sheets": len(wb.worksheets), "rows_sampled": row_count},
    )


def _extract_csv(content: bytes) -> Extraction:
    import csv
    import io

    reader = csv.reader(io.StringIO(content.decode("utf-8", errors="replace")))
    rows = list(reader)
    sample = rows[:20]
    text = "\n".join(", ".join(row) for row in sample)
    return Extraction(text=text[:_MAX_EXTRACT_CHARS], metadata={"rows": len(rows)})


def _extract_pptx(content: bytes) -> Extraction:
    import io

    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slide_texts = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        slide_texts.append(" / ".join(texts))
    text = "\n".join(slide_texts)
    return Extraction(text=text[:_MAX_EXTRACT_CHARS], metadata={"slides": len(prs.slides)})


def _extract_text(content: bytes) -> Extraction:
    text = content.decode("utf-8", errors="replace")
    return Extraction(text=text[:_MAX_EXTRACT_CHARS], metadata={"chars": len(text), "words": len(text.split())})


def _extract_audio(content: bytes, suffix: str) -> Extraction:
    from confector_ai_tools import transcribe

    result = transcribe(content, suffix)
    if result.error:
        return Extraction(text=None, metadata={"error": result.error})
    metadata = {"language": result.language, "duration_seconds": result.duration_seconds, "segments": result.segments}
    return Extraction(text=(result.text or "")[:_MAX_EXTRACT_CHARS] or None, metadata=metadata)
